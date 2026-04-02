
import os
from pptx import Presentation
from pptx.util import Inches

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Get blank layout
blank_layout = prs.slide_layouts[6]

# Add each image as a slide
image_dir = '/Users/ericcoffie/Market Assasin/market-assassin/presentations/slide-images'
for i in range(1, 98 + 1):
    img_path = os.path.join(image_dir, f'slide-{i:03d}.png')
    if os.path.exists(img_path):
        slide = prs.slides.add_slide(blank_layout)
        # Add image to fill the slide
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))
        print(f'Added slide {i}')

prs.save('/Users/ericcoffie/Market Assasin/market-assassin/presentations/JTED-2026-Final.pptx')
print(f'\nSaved JTED-2026-Final.pptx with {len(prs.slides)} slides')
