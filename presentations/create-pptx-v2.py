
import os
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]
image_dir = '/Users/ericcoffie/Market Assasin/market-assassin/presentations/slide-images-v2'

for i in range(1, 98 + 1):
    img_path = os.path.join(image_dir, f'slide-{i:03d}.png')
    if os.path.exists(img_path):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))

output_path = '/Users/ericcoffie/Market Assasin/market-assassin/presentations/JTED-2026-Final.pptx'
prs.save(output_path)
print(f'Saved {len(prs.slides)} slides to {output_path}')
