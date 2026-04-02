#!/usr/bin/env python3
"""
Convert JTED-2026-Revised.html to PowerPoint with formatting
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup
import re

# Color palette from the HTML
COLORS = {
    'bg_dark': RGBColor(15, 23, 42),      # #0f172a
    'bg_medium': RGBColor(30, 41, 59),     # #1e293b
    'green': RGBColor(16, 185, 129),       # #10b981
    'green_dark': RGBColor(5, 150, 105),   # #059669
    'purple': RGBColor(167, 139, 250),     # #a78bfa
    'purple_dark': RGBColor(124, 58, 237), # #7c3aed
    'warning': RGBColor(245, 158, 11),     # #f59e0b
    'alert': RGBColor(239, 68, 68),        # #ef4444
    'navy': RGBColor(30, 58, 138),         # #1e3a8a
    'white': RGBColor(255, 255, 255),
    'gray': RGBColor(148, 163, 184),       # #94a3b8
    'text': RGBColor(226, 232, 240),       # #e2e8f0
}

def clean_text(text):
    """Clean up text from HTML"""
    if not text:
        return ""
    text = re.sub(r'\s+', ' ', text)
    text = text.strip()
    return text

def get_slide_type(slide_div):
    """Determine slide type from classes"""
    classes = slide_div.get('class', [])
    if 'section-slide' in classes:
        if 'purple' in classes:
            return 'section-purple'
        elif 'navy' in classes:
            return 'section-navy'
        return 'section-green'
    if 'cta-slide' in classes:
        return 'cta'
    return 'content'

def extract_slide_content(slide_div):
    """Extract content from a slide div"""
    content = {
        'type': get_slide_type(slide_div),
        'title': '',
        'subtitle': '',
        'body': [],
        'stats': [],
        'quote': '',
        'quote_source': '',
        'icon': '',
        'big_stat': '',
        'big_stat_label': '',
        'insight': '',
        'slide_number': '',
    }

    # Get slide number
    slide_num = slide_div.find(class_='slide-number')
    if slide_num:
        content['slide_number'] = clean_text(slide_num.get_text())

    # Get icon (for section slides)
    icon = slide_div.find(class_='icon')
    if icon and not icon.find_parent(class_=['icon-box', 'stat-box', 'vs-box']):
        content['icon'] = clean_text(icon.get_text())

    # Get h1 title
    h1 = slide_div.find('h1')
    if h1:
        content['title'] = clean_text(h1.get_text())

    # Get h2 title
    h2 = slide_div.find('h2')
    if h2:
        if not content['title']:
            content['title'] = clean_text(h2.get_text())
        else:
            content['subtitle'] = clean_text(h2.get_text())

    # Get subtitle
    subtitle = slide_div.find(class_='subtitle')
    if subtitle:
        content['subtitle'] = clean_text(subtitle.get_text())

    # Get big stat
    big_stat = slide_div.find(class_='big-stat')
    if big_stat:
        number = big_stat.find(class_='number')
        label = big_stat.find(class_='label')
        if number:
            content['big_stat'] = clean_text(number.get_text())
        if label:
            content['big_stat_label'] = clean_text(label.get_text())

    # Get quote
    quote_box = slide_div.find(class_='quote-box')
    if quote_box:
        source = quote_box.find(class_='source')
        if source:
            content['quote_source'] = clean_text(source.get_text())
            source.decompose()
        content['quote'] = clean_text(quote_box.get_text())

    # Get stat boxes
    stat_boxes = slide_div.find_all(class_='stat-box')
    for box in stat_boxes:
        number = box.find(class_='number')
        label = box.find(class_='label')
        if number and label:
            content['stats'].append({
                'number': clean_text(number.get_text()),
                'label': clean_text(label.get_text())
            })

    # Get icon boxes as body items
    icon_boxes = slide_div.find_all(class_='icon-box')
    for box in icon_boxes:
        text = box.find(class_='text')
        if text:
            content['body'].append(clean_text(text.get_text()))

    # Get insight box
    insight = slide_div.find(class_='insight-box')
    if insight:
        text = insight.find(class_='text')
        if text:
            content['insight'] = clean_text(text.get_text())

    # Get table data
    table = slide_div.find('table')
    if table:
        rows = table.find_all('tr')
        for row in rows:
            cells = row.find_all(['th', 'td'])
            if cells:
                row_text = ' | '.join(clean_text(cell.get_text()) for cell in cells)
                content['body'].append(row_text)

    # Get action items
    actions = slide_div.find_all(class_='action-item')
    for action in actions:
        text = action.find(class_='text')
        if text:
            content['body'].append(clean_text(text.get_text()))

    return content

def add_slide(prs, content):
    """Add a slide to the presentation"""

    # Use blank layout
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Set slide background
    background = slide.background
    fill = background.fill
    fill.solid()

    if content['type'] == 'section-green':
        fill.fore_color.rgb = COLORS['green_dark']
    elif content['type'] == 'section-purple':
        fill.fore_color.rgb = COLORS['purple_dark']
    elif content['type'] == 'section-navy':
        fill.fore_color.rgb = COLORS['navy']
    elif content['type'] == 'cta':
        fill.fore_color.rgb = COLORS['purple_dark']
    else:
        fill.fore_color.rgb = COLORS['bg_dark']

    # Add title
    if content['title']:
        left = Inches(0.5)
        top = Inches(0.8)
        width = Inches(9)
        height = Inches(1.5)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = content['title']
        p.font.size = Pt(40) if content['type'].startswith('section') else Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER if content['type'].startswith('section') else PP_ALIGN.LEFT

    # Add icon for section slides
    if content['icon'] and content['type'].startswith('section'):
        left = Inches(4.5)
        top = Inches(2)
        width = Inches(1)
        height = Inches(1)

        icon_box = slide.shapes.add_textbox(left, top, width, height)
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = content['icon']
        p.font.size = Pt(60)
        p.alignment = PP_ALIGN.CENTER

    # Add subtitle
    if content['subtitle']:
        left = Inches(0.5)
        top = Inches(2.3) if content['type'].startswith('section') else Inches(1.8)
        width = Inches(9)
        height = Inches(0.8)

        sub_box = slide.shapes.add_textbox(left, top, width, height)
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = content['subtitle']
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['gray']
        p.alignment = PP_ALIGN.CENTER if content['type'].startswith('section') else PP_ALIGN.LEFT

    # Add big stat
    if content['big_stat']:
        left = Inches(0.5)
        top = Inches(2.5)
        width = Inches(9)
        height = Inches(2)

        stat_box = slide.shapes.add_textbox(left, top, width, height)
        tf = stat_box.text_frame

        p = tf.paragraphs[0]
        p.text = content['big_stat']
        p.font.size = Pt(96)
        p.font.bold = True
        p.font.color.rgb = COLORS['green']
        p.alignment = PP_ALIGN.CENTER

        if content['big_stat_label']:
            p2 = tf.add_paragraph()
            p2.text = content['big_stat_label']
            p2.font.size = Pt(24)
            p2.font.color.rgb = COLORS['gray']
            p2.alignment = PP_ALIGN.CENTER

    # Add quote
    if content['quote']:
        left = Inches(0.75)
        top = Inches(2.5)
        width = Inches(8.5)
        height = Inches(3)

        quote_box = slide.shapes.add_textbox(left, top, width, height)
        tf = quote_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f'"{content["quote"]}"'
        p.font.size = Pt(24)
        p.font.italic = True
        p.font.color.rgb = COLORS['text']
        p.alignment = PP_ALIGN.CENTER

        if content['quote_source']:
            p2 = tf.add_paragraph()
            p2.text = content['quote_source']
            p2.font.size = Pt(16)
            p2.font.color.rgb = COLORS['green']
            p2.alignment = PP_ALIGN.CENTER

    # Add stats
    if content['stats']:
        num_stats = len(content['stats'])
        stat_width = Inches(2.5)
        gap = Inches(0.3)
        total_width = num_stats * stat_width + (num_stats - 1) * gap
        start_left = (Inches(10) - total_width) / 2

        for i, stat in enumerate(content['stats'][:4]):  # Max 4 stats
            left = start_left + i * (stat_width + gap)
            top = Inches(3)

            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, stat_width, Inches(1.8)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLORS['green_dark']
            shape.line.fill.background()

            # Add text to shape
            tf = shape.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.text = stat['number']
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = COLORS['white']
            p.alignment = PP_ALIGN.CENTER

            p2 = tf.add_paragraph()
            p2.text = stat['label']
            p2.font.size = Pt(14)
            p2.font.color.rgb = COLORS['white']
            p2.alignment = PP_ALIGN.CENTER

    # Add body content
    if content['body'] and not content['stats'] and not content['quote'] and not content['big_stat']:
        left = Inches(0.5)
        top = Inches(2.5)
        width = Inches(9)
        height = Inches(4)

        body_box = slide.shapes.add_textbox(left, top, width, height)
        tf = body_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content['body'][:10]):  # Max 10 items
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['text']
            p.space_after = Pt(8)

    # Add insight box
    if content['insight']:
        left = Inches(0.5)
        top = Inches(6)
        width = Inches(9)
        height = Inches(0.8)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS['green_dark']
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Eric's Take: {content['insight']}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.LEFT

    # Add slide number
    if content['slide_number']:
        left = Inches(8.5)
        top = Inches(6.8)
        width = Inches(1.2)
        height = Inches(0.3)

        num_box = slide.shapes.add_textbox(left, top, width, height)
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = content['slide_number']
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['gray']
        p.alignment = PP_ALIGN.RIGHT

def main():
    # Read HTML file
    with open('JTED-2026-Revised.html', 'r', encoding='utf-8') as f:
        html = f.read()

    soup = BeautifulSoup(html, 'html.parser')

    # Find all slides
    slides = soup.find_all('div', class_='slide')
    print(f"Found {len(slides)} slides")

    # Create presentation (16:9 aspect ratio)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Process each slide
    for i, slide_div in enumerate(slides):
        content = extract_slide_content(slide_div)
        print(f"Processing slide {i+1}: {content['title'][:50] if content['title'] else 'No title'}...")
        add_slide(prs, content)

    # Save
    output_file = 'JTED-2026-Revised.pptx'
    prs.save(output_file)
    print(f"\nSaved: {output_file}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == '__main__':
    main()
